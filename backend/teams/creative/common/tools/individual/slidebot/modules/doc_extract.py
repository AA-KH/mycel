"""
Document Extraction Module - Supports PDF, Word, PPT, Excel, TXT etc.
"""

from pathlib import Path

# Document processing libraries (optional import)
try:
    import pymupdf as fitz  # PyMuPDF for PDF
except ImportError:
    fitz = None
    print("Warning: PyMuPDF is not installed. PDF text extraction is unavailable.")

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None
    print("Warning: python-docx is not installed. Word document processing is unavailable.")

try:
    from pptx import Presentation
except ImportError:
    Presentation = None
    print("Warning: python-pptx is not installed. PPT document processing is unavailable.")

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None
    print("Warning: openpyxl is not installed. Excel processing is unavailable.")


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF"""
    if fitz is None:
        return ""
    try:
        doc = fitz.open(file_path)
        text_parts = []
        for page in doc:
            text_parts.append(page.get_text())
        doc.close()
        return "\n".join(text_parts)
    except Exception as e:
        print(f"PDF extraction error: {e}")
        return ""


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from Word document"""
    if DocxDocument is None:
        return ""
    try:
        doc = DocxDocument(file_path)
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        # Extract text from tables as well
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_parts.append(" | ".join(row_text))
        return "\n".join(text_parts)
    except Exception as e:
        print(f"Word extraction error: {e}")
        return ""


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPT"""
    if Presentation is None:
        return ""
    try:
        prs = Presentation(file_path)
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = [f"[Slide {slide_num}]"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
            if len(slide_texts) > 1:
                text_parts.append("\n".join(slide_texts))
        return "\n\n".join(text_parts)
    except Exception as e:
        print(f"PPT extraction error: {e}")
        return ""


def extract_text_from_xlsx(file_path: str) -> str:
    """Extract text from Excel"""
    if load_workbook is None:
        return ""
    try:
        wb = load_workbook(file_path, data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_texts = [f"[Sheet: {sheet_name}]"]
            for row in sheet.iter_rows(values_only=True):
                row_values = [str(cell) for cell in row if cell is not None]
                if row_values:
                    sheet_texts.append(" | ".join(row_values))
            if len(sheet_texts) > 1:
                text_parts.append("\n".join(sheet_texts))
        wb.close()
        return "\n\n".join(text_parts)
    except Exception as e:
        print(f"Excel extraction error: {e}")
        return ""


def extract_text_from_document(file_path: str, filename: str) -> str:
    """Extract text based on file type"""
    ext = Path(filename).suffix.lower()
    
    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext in ['.docx', '.doc']:
        return extract_text_from_docx(file_path)
    elif ext in ['.pptx', '.ppt']:
        return extract_text_from_pptx(file_path)
    elif ext in ['.xlsx', '.xls']:
        return extract_text_from_xlsx(file_path)
    elif ext == '.txt':
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            print(f"TXT read error: {e}")
            return ""
    else:
        return ""


def extract_table_from_file(file_path: str, filename: str) -> str:
    """
    Extract table content from Excel or CSV file, format as text suitable for AI comprehension
    Limit to maximum 3000 characters to avoid overly long prompts
    """
    ext = Path(filename).suffix.lower()
    MAX_TABLE_LENGTH = 3000
    
    if ext == '.csv':
        try:
            import csv
            rows = []
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f)
                for row in reader:
                    if any(cell.strip() for cell in row):
                        rows.append(" | ".join(row))
                    if len("\n".join(rows)) > MAX_TABLE_LENGTH:
                        break
            result = "\n".join(rows)
            if len(result) > MAX_TABLE_LENGTH:
                result = result[:MAX_TABLE_LENGTH] + "\n...(Data too long, truncated)"
            return result
        except Exception as e:
            print(f"CSV read error: {e}")
            return ""
    
    elif ext in ['.xlsx', '.xls']:
        if load_workbook is None:
            return ""
        try:
            wb = load_workbook(file_path, data_only=True)
            all_tables = []
            total_length = 0
            
            for sheet_name in wb.sheetnames:
                if total_length > MAX_TABLE_LENGTH:
                    break
                    
                sheet = wb[sheet_name]
                table_rows = []
                
                for row in sheet.iter_rows(values_only=True):
                    row_values = []
                    for cell in row:
                        if cell is not None:
                            row_values.append(str(cell))
                        else:
                            row_values.append("")
                    # Only add non-empty rows
                    if any(v.strip() for v in row_values):
                        table_rows.append(" | ".join(row_values))
                    
                    # Check length
                    if len("\n".join(table_rows)) > MAX_TABLE_LENGTH:
                        break
                
                if table_rows:
                    if len(wb.sheetnames) > 1:
                        all_tables.append(f"[{sheet_name}]\n" + "\n".join(table_rows))
                    else:
                        all_tables.append("\n".join(table_rows))
                    total_length = len("\n\n".join(all_tables))
            
            wb.close()
            result = "\n\n".join(all_tables)
            if len(result) > MAX_TABLE_LENGTH:
                result = result[:MAX_TABLE_LENGTH] + "\n...(Data too long, truncated)"
            return result
        except Exception as e:
            print(f"Excel table read error: {e}")
            return ""
    
    return ""
